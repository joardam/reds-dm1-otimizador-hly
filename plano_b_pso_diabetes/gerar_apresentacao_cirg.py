# -*- coding: utf-8 -*-
"""Apresentação Equipe 4 — Plano B, no template institucional CIRG/POLI-UPE (Prof. Fernando Buarque)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ---------- paleta do template ----------
NAVY  = RGBColor(0x20, 0x38, 0x64)   # títulos / texto institucional
STEEL = RGBColor(0x41, 0x71, 0x9C)   # filete do rodapé / acento
INK   = RGBColor(0x40, 0x40, 0x40)   # corpo
GRAY  = RGBColor(0x59, 0x59, 0x59)   # secundário
HAIR  = RGBColor(0xD3, 0xD9, 0xE0)
LIGHT = RGBColor(0xF1, 0xF4, 0xF8)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
ONNAVY= RGBColor(0xCD, 0xD7, 0xE3)
GREEN = RGBColor(0x2E, 0x6B, 0x4F)
AMBER = RGBColor(0xB0, 0x7A, 0x2A)
GREYS = RGBColor(0x9A, 0x9A, 0x9A)

SERIF = "Palatino Linotype"   # fonte exata do template do professor
SANSB = "Arial"               # marcadores de bullet (igual ao template)
ARROW = "→"

BASE = "/home/user/reds-dm1-otimizador-hly/plano_b_pso_diabetes"
A = BASE + "/assets_template"
HDR = A + "/header_right.png"
FTR = A + "/footer_left.png"
CONV_PNG = BASE + "/resultados_convergencia_pso.png"
IMP_PNG  = BASE + "/resultados_importancia.png"
META_PNG = BASE + "/meta_fss_pso/resultados_comparacao_convergencia.png"
OUT = BASE + "/Apresentacao_Equipe4_PlanoB_CIRG.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
ML = Inches(0.55)
CW = SW - Inches(1.1)

def _f(run, size=18, color=INK, bold=False, italic=False, font=SERIF):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.name = font; run.font.color.rgb = color

def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf

def para(tf, text, size=18, color=INK, bold=False, italic=False, font=SERIF,
         align=PP_ALIGN.LEFT, space_after=6, space_before=0, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    r = p.add_run(); r.text = text; _f(r, size, color, bold, italic, font)
    return p

def rect(slide, l, t, w, h, fill=LIGHT, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h); sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp

def bullets(tf, items, size=16, gap=8, color=INK, lead=NAVY, first0=True):
    for i, it in enumerate(items):
        txt, bold = (it if isinstance(it, tuple) else (it, False))
        p = tf.paragraphs[0] if (i == 0 and first0) else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        r0 = p.add_run(); r0.text = "•  "; _f(r0, size, lead, True, False, SANSB)
        r = p.add_run(); r.text = txt; _f(r, size, color, bold, False, SERIF)

def master(slide, num):
    rect(slide, 0, 0, SW, SH, fill=PAPER)
    hw = Inches(4.55)
    slide.shapes.add_picture(HDR, int(SW - hw), Inches(0.12), width=hw)
    rect(slide, Inches(0.3), Inches(6.82), SW - Inches(0.6), Pt(1.4), fill=STEEL)
    slide.shapes.add_picture(FTR, Inches(0.32), Inches(6.95), height=Inches(0.5))
    _, ctf = box(slide, Inches(1.85), Inches(6.97), Inches(9.0), Inches(0.5))
    para(ctf, "COURSE: Computação Natural", size=11, color=NAVY, first=True, space_after=1)
    para(ctf, "CLASS: Equipe 4  -  TOPIC: Otimizador Bioinspirado (Plano B)", size=11, color=NAVY, space_after=0)
    _, ntf = box(slide, SW - Inches(0.95), Inches(7.0), Inches(0.6), Inches(0.3))
    para(ntf, str(num), size=12, color=NAVY, align=PP_ALIGN.RIGHT, first=True)

def titled(title, num):
    s = prs.slides.add_slide(BLANK); master(s, num)
    _, ttf = box(s, ML, Inches(0.42), Inches(8.0), Inches(0.95))
    para(ttf, title, size=27, color=NAVY, bold=True, first=True)
    return s

def badge(s, x, y, w, text, col):
    rect(s, x, y, w, Inches(0.34), fill=None, line=col, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _, tf = box(s, x, y+Inches(0.04), w, Inches(0.3))
    para(tf, text.upper(), size=10.5, color=col, bold=True, align=PP_ALIGN.CENTER, first=True)

# ============================================================ TÍTULO (1)
s = prs.slides.add_slide(BLANK); master(s, 1)
_, tg = box(s, Inches(5.0), Inches(0.92), SW - Inches(5.3), Inches(0.35))
para(tg, "At CIRG of POLI/UPE Nature insights nurture our creativity for real problem solving",
     size=11, color=NAVY, bold=True, italic=True, align=PP_ALIGN.RIGHT, first=True)
_, tt = box(s, Inches(1.0), Inches(2.25), SW - Inches(2.0), Inches(1.7))
para(tt, "Projeto REDS-DM1", size=38, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=4)
para(tt, "Otimizador Bioinspirado para Gestão de Saúde em Diabetes  ·  Plano B",
     size=21, color=NAVY, align=PP_ALIGN.CENTER, space_after=0)
_, au = box(s, Inches(1.0), Inches(3.95), SW - Inches(2.0), Inches(1.4))
para(au, "Pedro Henrique França  ·  Silas Nunes  ·  Jéssica Maria  ·  José Mario",
     size=19, color=GRAY, align=PP_ALIGN.CENTER, first=True, space_after=6)
para(au, "Prof. Dr. Fernando Buarque de Lima Neto", size=19, color=GRAY, align=PP_ALIGN.CENTER, space_after=4)
para(au, "Mentoria: Rodrigo Amaro  ·  Allan Miller", size=16, color=GRAY, align=PP_ALIGN.CENTER, space_after=0)
_, af = box(s, Inches(1.0), Inches(5.5), SW - Inches(2.0), Inches(1.2))
para(af, "Universidade de Pernambuco – UPE", size=16, color=GRAY, align=PP_ALIGN.CENTER, first=True, space_after=2)
para(af, "Escola Politécnica de Pernambuco – POLI", size=16, color=GRAY, align=PP_ALIGN.CENTER, space_after=2)
para(af, "Computational Intelligent Research Group – CIRG", size=16, color=GRAY, align=PP_ALIGN.CENTER, space_after=0)

# ============================================================ 2 AGENDA
s = titled("Agenda", 2)
items = [
    ("01", "De onde viemos", "Trajetória em 3 fases e redução de escopo"),
    ("02", "O Plano B e seu objetivo", "Objetivo adaptado e a variável Wellness"),
    ("03", "Como funciona", "Passo a passo e por que PSO (enxame)"),
    ("04", "Validação e ganho", "O ótimo conhecido, o perfil ideal e o Meta-FSS"),
    ("05", "Fechamento", "BFSS de relance, limitações e próximos passos"),
]
top = Inches(1.55); rowh = Inches(0.86); gap = Inches(0.07)
for i,(n,t,d) in enumerate(items):
    y = top + i*(rowh+gap)
    rect(s, ML, y, CW, rowh, fill=(LIGHT if i in (1,2,3) else PAPER), line=HAIR, line_w=0.75)
    rect(s, ML, y, Inches(0.07), rowh, fill=STEEL)
    _, ntf = box(s, ML+Inches(0.28), y+Inches(0.16), Inches(0.9), Inches(0.6))
    para(ntf, n, size=22, color=STEEL, bold=True, first=True)
    _, ttf = box(s, ML+Inches(1.3), y+Inches(0.12), CW-Inches(1.5), Inches(0.7))
    para(ttf, t, size=17, color=NAVY, bold=True, first=True, space_after=1)
    para(ttf, d, size=12.5, color=GRAY, space_after=0)

# ============================================================ 3 SÍNTESE
s = titled("Síntese executiva", 3)
colw = (CW - Inches(0.6)) / 3
cols = [
    ("Contexto", NAVY, ["Disciplina de Computação Natural.","Sistema REDS-PE (Rede de Atenção à Saúde).",
        "Alvo: pacientes com diabetes em Pernambuco.","Financiamento: SUS / políticas públicas."]),
    ("Proposta", NAVY, ["Usar Inteligência de Enxames (PSO).","Achar o perfil do 'indivíduo ideal' entre diabéticos.",
        "Horizonte: otimizar a alocação de recursos maximizando HLY."]),
    ("Resultado", GREEN, ["PSO implementado do zero (NumPy).","Validado contra a solução analítica: cosseno = 1,000.",
        "Perfil clínico interpretável extraído dos pesos."]),
]
ctop = Inches(1.7)
for i,(h,bc,its) in enumerate(cols):
    x = ML + i*(colw+Inches(0.3))
    rect(s, x, ctop, colw, Inches(4.4), fill=LIGHT, line=HAIR, line_w=0.75)
    rect(s, x, ctop, colw, Inches(0.1), fill=bc)
    _, tf = box(s, x+Inches(0.28), ctop+Inches(0.35), colw-Inches(0.56), Inches(3.9))
    para(tf, h, size=18, color=bc, bold=True, first=True, space_after=10)
    bullets(tf, its, size=14, gap=10, lead=bc, first0=False)

# ============================================================ 4 TRAJETÓRIA
s = titled("Trajetória em três fases", 4)
_, sub = box(s, ML, Inches(1.32), CW, Inches(0.35))
para(sub, "A redução de escopo foi decisão técnica, não recuo.", size=14, color=GRAY, italic=True, first=True)
phases = [
    ("FASE 1","130-US + CTGAN","Abandonada",GREYS,["População errada (97% DM2)","Sem eixo temporal","Sem gabarito de validação"]),
    ("FASE 2","Marcadores plantados","Pausada",AMBER,["Base sintética com gabarito","Prova algorítmica do otimizador","BFSS validado aqui"]),
    ("FASE 3","Plano B — PSO / coorte bimodal","Entregue",GREEN,["Foco da apresentação de hoje","PSO validado (cosseno 1,000)","Perfil ideal interpretável"]),
]
cw3 = (CW - Inches(1.0))/3; ytop = Inches(2.05)
rect(s, ML+Inches(0.3), ytop, CW-Inches(0.6), Pt(2), fill=HAIR)
for i,(ph,name,st,col,pts) in enumerate(phases):
    x = ML + i*(cw3+Inches(0.5))
    rect(s, x+cw3/2-Inches(0.1), ytop-Inches(0.1), Inches(0.2), Inches(0.2), fill=col, shape=MSO_SHAPE.OVAL)
    ct = ytop + Inches(0.4)
    rect(s, x, ct, cw3, Inches(3.55), fill=PAPER, line=HAIR, line_w=0.75)
    rect(s, x, ct, cw3, Inches(0.09), fill=col)
    _, tf = box(s, x+Inches(0.25), ct+Inches(0.28), cw3-Inches(0.5), Inches(3.1))
    para(tf, ph, size=11, color=GRAY, bold=True, first=True, space_after=2)
    para(tf, name, size=16, color=NAVY, bold=True, space_after=5)
    para(tf, st.upper(), size=11, color=col, bold=True, space_after=9)
    bullets(tf, pts, size=12.5, gap=7, lead=col, first0=False)

# ============================================================ 5 REDUÇÃO DE ESCOPO
s = titled("Redução de escopo — status de cada componente", 5)
rows = [("130-US + CTGAN","Abandonado",GREYS),("PSO sobre coorte bimodal (Plano B)","Entregue",GREEN),
    ("Meta-FSS (hiperparâmetros do PSO)","Entregue",GREEN),("Visualização HTML do PSO (3 abas)","Entregue",GREEN),
    ("Banco sintético de marcadores","Concluído",GREEN),("BFSS (seleção de variáveis)","Validado",GREEN),
    ("MOPSO + NSGA-II (frente de Pareto)","Planejado",AMBER)]
ytop = Inches(1.55); rh = Inches(0.6); gap = Inches(0.07)
for i,(comp,st,col) in enumerate(rows):
    y = ytop + i*(rh+gap); bold = "Plano B" in comp
    rect(s, ML, y, CW, rh, fill=(LIGHT if bold else PAPER), line=HAIR, line_w=0.75)
    rect(s, ML, y, Inches(0.09), rh, fill=col)
    _, tf = box(s, ML+Inches(0.32), y+Inches(0.13), CW-Inches(3.0), Inches(0.4))
    para(tf, comp, size=15.5, color=(NAVY if bold else INK), bold=bold, first=True)
    badge(s, ML+CW-Inches(2.1), y+Inches(0.13), Inches(1.9), st, col)

# ============================================================ 6 OBJETIVO
s = titled(f"Objetivo: original  {ARROW}  adaptado", 6)
colw = (CW - Inches(1.1))/2; ytop = Inches(1.6); ch = Inches(3.75)
rect(s, ML, ytop, colw, ch, fill=PAPER, line=HAIR, line_w=0.75)
rect(s, ML, ytop, colw, Inches(0.1), fill=GREYS)
_, tf = box(s, ML+Inches(0.3), ytop+Inches(0.32), colw-Inches(0.6), ch-Inches(0.6))
para(tf, "ORIGINAL — PLANO A", size=12, color=GRAY, bold=True, first=True, space_after=8)
para(tf, "Otimizador multiobjetivo completo", size=17, color=NAVY, bold=True, space_after=10)
bullets(tf, ["Três objetivos concorrentes (frente de Pareto):","Maximizar HLY (Anos de Vida Saudáveis).",
    "Minimizar o custo da política pública.","Maximizar a equidade regional entre municípios.",
    "Baselines: MOPSO + NSGA-II."], size=13.5, gap=8, lead=GREYS, first0=False)
_, atf = box(s, ML+colw+Inches(0.15), ytop+ch/2-Inches(0.45), Inches(0.8), Inches(0.8))
para(atf, ARROW, size=40, color=STEEL, bold=True, align=PP_ALIGN.CENTER, first=True)
x2 = ML + colw + Inches(1.1)
rect(s, x2, ytop, colw, ch, fill=LIGHT, line=NAVY, line_w=1.25)
rect(s, x2, ytop, colw, Inches(0.1), fill=NAVY)
_, t2 = box(s, x2+Inches(0.3), ytop+Inches(0.32), colw-Inches(0.6), ch-Inches(0.6))
para(t2, "ADAPTADO — PLANO B  (entregável real)", size=12, color=NAVY, bold=True, first=True, space_after=8)
para(t2, "PSO acha o perfil do indivíduo ideal", size=17, color=NAVY, bold=True, space_after=10)
bullets(t2, ["PSO otimiza os pesos de uma regressão logística.","Variável-alvo: Wellness (proxy de HLY).",
    "Validar contra a solução analítica do scikit-learn.","Extrair o perfil clínico interpretável dos pesos."],
    size=13.5, gap=8, lead=NAVY, first0=False)
_, ntf = box(s, ML, ytop+ch+Inches(0.16), CW, Inches(0.5))
para(ntf, "Honestidade metodológica: a base é testbed; o entregável é o otimizador. Coorte única (China, 2012), que não distingue T1D/T2D.",
     size=12, color=GRAY, italic=True, first=True)

# ============================================================ 7 WELLNESS
s = titled("A variável dependente: Wellness", 7)
lw = CW*0.56
_, tf = box(s, ML, Inches(1.6), lw, Inches(4.0))
para(tf, "Por que não um rótulo clínico direto?", size=17, color=NAVY, bold=True, first=True, space_after=10)
bullets(tf, ["A base não possui rótulo 'saudável vs. doente'.","É um snapshot transversal, sem trajetória longitudinal.",
    "Wellness é construído a partir de variáveis objetivas da coorte."], size=14.5, gap=10, first0=False)
para(tf, "Conexão com HLY", size=15, color=NAVY, bold=True, space_after=6, space_before=8)
para(tf, "Ausência de complicações + controle metabólico = pilares que sustentam os Anos de Vida Saudáveis.",
     size=14, color=INK, space_after=0)
fy = Inches(4.6)
rect(s, ML, fy, lw, Inches(0.95), fill=LIGHT, line=HAIR, line_w=0.75)
_, ftf = box(s, ML+Inches(0.3), fy+Inches(0.13), lw-Inches(0.6), Inches(0.75))
para(ftf, "Wellness  =  terço superior de", size=13, color=GRAY, first=True, space_after=2)
para(ftf, "z(idade) − z(HbA1c) − z(n.º de comorbidades)", size=18, color=NAVY, bold=True, space_after=0)
rx = ML + lw + Inches(0.5); rw = CW - lw - Inches(0.5); gy = Inches(1.6)
rect(s, rx, gy, rw, Inches(3.0), fill=PAPER, line=HAIR, line_w=0.75)
rect(s, rx, gy, rw, Inches(0.55), fill=NAVY)
_, htf = box(s, rx, gy+Inches(0.12), rw, Inches(0.4))
para(htf, "Grupos  (filtro DM == 1)", size=14, color=PAPER, bold=True, align=PP_ALIGN.CENTER, first=True)
grp = [("Indivíduos ideais","119","terço superior"),("Demais diabéticos","238",""),("Total","357","")]
yy = gy + Inches(0.72)
for i,(g,n,note) in enumerate(grp):
    if i==2: rect(s, rx+Inches(0.2), yy, rw-Inches(0.4), Pt(1), fill=HAIR); yy += Inches(0.08)
    _, rtf = box(s, rx+Inches(0.3), yy, rw-Inches(1.4), Inches(0.5))
    para(rtf, g, size=15, color=INK, bold=(i in (0,2)), first=True, space_after=0)
    if note: para(rtf, note, size=10.5, color=GRAY, italic=True, space_after=0)
    _, vtf = box(s, rx+rw-Inches(1.2), yy-Inches(0.02), Inches(1.0), Inches(0.5))
    para(vtf, n, size=20, color=NAVY, bold=True, align=PP_ALIGN.RIGHT, first=True)
    yy += Inches(0.72)

# ============================================================ 8 PASSO A PASSO
s = titled("Passo a passo da solução (Plano B)", 8)
_, sub = box(s, ML, Inches(1.3), CW, Inches(0.35))
para(sub, "Da base bruta ao perfil do indivíduo ideal — 9 etapas reprodutíveis.", size=13.5, color=GRAY, italic=True, first=True)
groups = [
    ("PREPARAÇÃO", STEEL, [("1","Carregar & filtrar","CSV 5.922×190 → DM==1 → 357"),
        ("2","Construir Wellness","rótulo 0/1"),("3","Selecionar 26 features","imputar mediana; z-score")]),
    ("OTIMIZAÇÃO", NAVY, [("4","Função de fitness","−(log-loss + L2)"),
        ("5","Executar PSO","27D; w=0,7; c1=c2=1,5"),("6","Validar","cosseno PSO × sklearn = 1,000")]),
    ("INTERPRETAÇÃO", GREEN, [("7","Meta-FSS","hiperparâm.; AUC cross-val"),
        ("8","Perfil do ideal","rank por |peso|; sinal = direção"),("9","Visualizar","HTML, 3 abas, enxame animado")]),
]
gw = (CW - Inches(0.8))/3; gy = Inches(1.85)
for i,(gh,col,steps) in enumerate(groups):
    x = ML + i*(gw+Inches(0.4))
    _, ghtf = box(s, x, gy, gw, Inches(0.35))
    para(ghtf, gh, size=12, color=col, bold=True, align=PP_ALIGN.CENTER, first=True)
    sy = gy + Inches(0.5)
    for (n,t,d) in steps:
        rect(s, x, sy, gw, Inches(1.12), fill=LIGHT, line=HAIR, line_w=0.75)
        rect(s, x, sy, Inches(0.1), Inches(1.12), fill=col)
        rect(s, x+Inches(0.28), sy+Inches(0.31), Inches(0.5), Inches(0.5), fill=col, shape=MSO_SHAPE.OVAL)
        _, ntf = box(s, x+Inches(0.28), sy+Inches(0.4), Inches(0.5), Inches(0.4))
        para(ntf, n, size=16, color=PAPER, bold=True, align=PP_ALIGN.CENTER, first=True)
        _, stf = box(s, x+Inches(0.95), sy+Inches(0.2), gw-Inches(1.1), Inches(0.8))
        para(stf, t, size=14.5, color=NAVY, bold=True, first=True, space_after=2)
        para(stf, d, size=11.5, color=GRAY, space_after=0)
        sy += Inches(1.22)

# ============================================================ 9 POR QUE PSO
s = titled("Por que PSO (enxame) e não algoritmo evolucionário", 9)
quad = [
    ("Terreno do problema", NAVY, "Os pesos vivem em espaço contínuo de alta dimensão (27D) e a perda é convexa. O PSO desliza no contínuo com momento e converge rápido para o ótimo. GAs exigem crossover/mutação, menos intuitivos para pesos reais."),
    ("Alinhamento com o enunciado", STEEL, "A disciplina avalia Inteligência de Enxames. PSO é o algoritmo canônico, e a família FSS é o fio condutor do projeto."),
    ("Interpretabilidade", STEEL, "O PSO entrega direto um vetor de pesos: magnitude = importância, sinal = direção clínica. Um genoma binário de GA obscureceria isso."),
    ("Validação elegante", GREEN, "Em problema convexo existe ótimo conhecido — dá para provar que o otimizador chegou nele (cosseno = 1,000). Demonstração limpa e didática."),
]
qw = (CW - Inches(0.5))/2; qh = Inches(2.2); qy0 = Inches(1.55)
for i,(h,col,body) in enumerate(quad):
    x = ML + (i%2)*(qw+Inches(0.5)); y = qy0 + (i//2)*(qh+Inches(0.3))
    rect(s, x, y, qw, qh, fill=LIGHT, line=HAIR, line_w=0.75)
    rect(s, x, y, Inches(0.1), qh, fill=col)
    _, tf = box(s, x+Inches(0.35), y+Inches(0.25), qw-Inches(0.65), qh-Inches(0.5))
    para(tf, h, size=16.5, color=col, bold=True, first=True, space_after=7)
    para(tf, body, size=13.5, color=INK, space_after=0)

# ============================================================ 10 VALIDAÇÃO
s = titled("O otimizador acha o ótimo conhecido", 10)
lw = CW*0.42
metrics = [("AUC treino","0,820",GRAY),("AUC teste","0,664",INK),("AUC sklearn (baseline)","0,664",INK)]
my = Inches(1.65)
for (lbl,val,col) in metrics:
    rect(s, ML, my, lw, Inches(0.68), fill=PAPER, line=HAIR, line_w=0.75)
    _, ltf = box(s, ML+Inches(0.25), my+Inches(0.16), lw-Inches(1.5), Inches(0.4))
    para(ltf, lbl, size=14, color=col, first=True)
    _, vtf = box(s, ML+lw-Inches(1.4), my+Inches(0.08), Inches(1.2), Inches(0.5))
    para(vtf, val, size=22, color=col, bold=True, align=PP_ALIGN.RIGHT, first=True)
    my += Inches(0.8)
cy = my + Inches(0.05)
rect(s, ML, cy, lw, Inches(1.3), fill=NAVY)
_, ctf = box(s, ML+Inches(0.25), cy+Inches(0.2), lw-Inches(0.5), Inches(1.0))
para(ctf, "Cosseno(pesos PSO, sklearn)", size=13, color=ONNAVY, bold=True, first=True, space_after=2)
para(ctf, "1,000", size=38, color=PAPER, bold=True, space_after=0)
rx = ML + lw + Inches(0.5); rw = CW - lw - Inches(0.5)
rect(s, rx, Inches(1.6), rw, Inches(3.5), fill=PAPER, line=HAIR, line_w=0.75)
s.shapes.add_picture(CONV_PNG, rx+Inches(0.15), Inches(1.75), width=rw-Inches(0.3))
_, ntf = box(s, ML, Inches(5.55), CW, Inches(0.95))
para(ntf, "O AUC mede a dificuldade do problema; o cosseno prova o otimizador. Para uma demonstração de validação, o PSO num problema convexo é limpo e didático: encontra EXATAMENTE a mesma solução que o otimizador analítico.",
     size=12.5, color=GRAY, italic=True, first=True)

# ============================================================ 11 GANHO
s = titled("O ganho: o perfil do indivíduo ideal", 11)
rw = CW*0.5; rx = ML + CW - rw
rect(s, rx, Inches(1.55), rw, Inches(4.7), fill=PAPER, line=HAIR, line_w=0.75)
pic = s.shapes.add_picture(IMP_PNG, rx, Inches(1.7), height=Inches(4.4)); pic.left = int(rx + (rw - pic.width)/2)
lw = CW - rw - Inches(0.5)
_, tf = box(s, ML, Inches(1.6), lw, Inches(2.2))
para(tf, "Top características (peso do PSO)", size=15, color=NAVY, bold=True, first=True, space_after=8)
bullets(tf, ["CP2h +0,338  ·  BMI −0,304  ·  BUN +0,240","waist1 −0,232  ·  ISIGutt −0,212  ·  WHR −0,193",
    "SCRE +0,162  ·  ALT −0,134  ·  FatherDM −0,119"], size=13.5, gap=7, first0=False)
py = Inches(3.35)
rect(s, ML, py, lw, Inches(1.55), fill=LIGHT, line=NAVY, line_w=1.0)
_, ptf = box(s, ML+Inches(0.25), py+Inches(0.18), lw-Inches(0.5), Inches(1.2))
para(ptf, "Perfil-síntese do diabético ideal", size=14, color=NAVY, bold=True, first=True, space_after=5)
para(ptf, "Mais magro, com menor resistência à insulina, função hepática preservada (ALT) e sem histórico familiar de DM.",
     size=13.5, color=INK, space_after=0)
_, mtf = box(s, ML, Inches(5.15), lw, Inches(1.1))
para(mtf, "Médias (ideais vs demais):", size=13, color=GRAY, bold=True, first=True, space_after=4)
para(mtf, "BMI 21,4 vs 24,8   ·   ISIGutt 57,3 vs 104,8   ·   ALT 33,4 vs 45,7", size=13, color=INK, space_after=0)

# ============================================================ 12 META-FSS
s = titled("Meta-FSS — o enxame que ajusta o enxame", 12)
_, sub = box(s, ML, Inches(1.3), CW, Inches(0.35))
para(sub, "FSS (Fish School Search) otimiza os hiperparâmetros do PSO — meta-otimização em dois níveis.", size=13.5, color=GRAY, italic=True, first=True)
flow = [("FSS — meta-nível","ajusta w, c1, c2, λ",NAVY),("PSO — nível 1","otimiza os pesos do modelo",STEEL),("Regressão logística","classifica o Wellness",GREEN)]
fbw = (CW - Inches(1.0))/3; fy = Inches(1.85)
for i,(h,d,col) in enumerate(flow):
    x = ML + i*(fbw+Inches(0.5))
    rect(s, x, fy, fbw, Inches(1.0), fill=LIGHT, line=col, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, fy, Inches(0.1), Inches(1.0), fill=col)
    _, tf = box(s, x+Inches(0.32), fy+Inches(0.2), fbw-Inches(0.5), Inches(0.7))
    para(tf, h, size=14.5, color=NAVY, bold=True, first=True, space_after=2)
    para(tf, d, size=12, color=GRAY, space_after=0)
    if i < 2:
        _, atf = box(s, x+fbw+Inches(0.04), fy+Inches(0.24), Inches(0.42), Inches(0.5))
        para(atf, ARROW, size=22, color=GRAY, bold=True, align=PP_ALIGN.CENTER, first=True)
lw = CW*0.46; ly = Inches(3.25)
_, tf = box(s, ML, ly, lw, Inches(2.3))
para(tf, "Hiperparâmetros encontrados pelo FSS", size=14.5, color=NAVY, bold=True, first=True, space_after=9)
bullets(tf, ["w = 0,67  ·  c1 = 0,74  ·  c2 = 0,59  ·  λ = 0,016",
    "AUC teste: 0,678 (meta)  vs  0,664 (padrão)  →  +0,014","Avaliação por CV-3 no treino (sem tocar no teste)."],
    size=13.5, gap=10, first0=False)
rw = CW - lw - Inches(0.5); rx = ML + lw + Inches(0.5)
rect(s, rx, Inches(3.15), rw, Inches(2.4), fill=PAPER, line=HAIR, line_w=0.75)
pic = s.shapes.add_picture(META_PNG, rx, Inches(3.3), height=Inches(2.15)); pic.left = int(rx + (rw - pic.width)/2)
cy = Inches(5.75)
rect(s, ML, cy, CW, Inches(0.82), fill=LIGHT, line=STEEL, line_w=1.0)
_, ctf = box(s, ML+Inches(0.3), cy+Inches(0.12), CW-Inches(0.6), Inches(0.62))
para(ctf, "Leitura honesta: o ganho (+0,014) está dentro do ruído da estimativa (IC ≈ 0,68 ± 0,15) e parte vem do afrouxamento do L2. O ponto não é o número — é mostrar a família FSS governando o PSO.",
     size=12, color=INK, italic=True, first=True)

# ============================================================ 13 BFSS
s = titled("BFSS, de relance — o plano principal segue vivo", 13)
_, tf = box(s, ML, Inches(1.7), CW*0.58, Inches(4.0))
para(tf, "Validação por marcadores plantados", size=17, color=NAVY, bold=True, first=True, space_after=10)
bullets(tf, ["Base sintética com gabarito (verdade-base plantada).",
    "Prova ALGORÍTMICA de que o otimizador recupera o sinal — não é alegação clínica.",
    "Mantém a família de enxame (FSS) como fio condutor do projeto."], size=14.5, gap=10, first0=False)
para(tf, "Próximo: MOPSO discreto + NSGA-II para a frente de Pareto.", size=14, color=GRAY, italic=True, space_after=0, space_before=6)
rx = ML + CW*0.62; rw = CW*0.38
rect(s, rx, Inches(1.7), rw, Inches(1.95), fill=LIGHT, line=HAIR, line_w=0.75)
rect(s, rx, Inches(1.7), rw, Inches(0.09), fill=GREEN)
_, stf = box(s, rx+Inches(0.3), Inches(1.95), rw-Inches(0.6), Inches(1.6))
para(stf, "STATUS", size=11, color=GRAY, bold=True, first=True, space_after=4)
para(stf, "BFSS validado", size=18, color=GREEN, bold=True, space_after=6)
para(stf, "Recuperou 8 de 9 marcadores plantados; o subconjunto selecionado eleva o R².", size=13, color=INK, space_after=0)

# ============================================================ 14 LIMITAÇÕES
s = titled("Limitações declaradas", 14)
lims = ["A base não distingue T1D/T2D → 'diabetes tipo não especificado'; é testbed, não evidência clínica.",
    "Coorte única e transversal (China, 2012) → fala de associação, não de causalidade.",
    "O rótulo 'ideal' é uma escolha de projeto defensável, não a única.",
    "As 26 preditoras foram curadas à mão → podem ter omitido pistas relevantes.",
    "Imputação pela mediana 'achata' a variação; o corte do terço (33%) é limiar arbitrário.",
    "L2 (λ=0,05) calibrado empiricamente, não por cross-validation formal."]
colw = (CW - Inches(0.5))/2; ytop = Inches(1.65)
for i,t in enumerate(lims):
    x = ML + (i%2)*(colw+Inches(0.5)); y = ytop + (i//2)*Inches(1.45)
    rect(s, x, y, colw, Inches(1.25), fill=LIGHT, line=HAIR, line_w=0.75)
    rect(s, x, y, Inches(0.09), Inches(1.25), fill=STEEL)
    _, tf = box(s, x+Inches(0.3), y+Inches(0.22), colw-Inches(0.55), Inches(0.95))
    para(tf, t, size=13.5, color=INK, first=True)

# ============================================================ 15 CONCLUSÃO
s = titled("Próximos passos e conclusão", 15)
stages = ["PSO (Plano B)","BFSS","Simulador HLY","MOPSO + NSGA-II"]
pw = (CW - Inches(1.2))/4; py = Inches(1.7)
for i,st in enumerate(stages):
    x = ML + i*(pw+Inches(0.4)); col = GREEN if i<2 else (AMBER if i==3 else STEEL)
    rect(s, x, py, pw, Inches(0.85), fill=LIGHT, line=col, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _, tf = box(s, x, py+Inches(0.26), pw, Inches(0.5))
    para(tf, st, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True)
    if i < 3:
        _, atf = box(s, x+pw+Inches(0.02), py+Inches(0.2), Inches(0.36), Inches(0.5))
        para(atf, ARROW, size=20, color=GRAY, bold=True, align=PP_ALIGN.CENTER, first=True)
_, tf = box(s, ML, Inches(3.05), CW, Inches(1.2))
para(tf, "Próximo passo — frente de Pareto", size=16, color=NAVY, bold=True, first=True, space_after=6)
para(tf, "MOPSO discreto + NSGA-II sobre ~540 políticas, com gabarito computável por enumeração (HLY × Custo × Equidade).",
     size=14.5, color=INK, space_after=0)
fy = Inches(4.35)
rect(s, ML, fy, CW, Inches(1.5), fill=NAVY)
_, ftf = box(s, ML+Inches(0.4), fy+Inches(0.3), CW-Inches(0.8), Inches(1.0))
para(ftf, "Conclusão", size=13, color=ONNAVY, bold=True, first=True, space_after=4)
para(ftf, "PSO validado (cosseno 1,000) + perfil clínico interpretável = o otimizador funciona. O Plano A continua como horizonte multiobjetivo.",
     size=17, color=PAPER, bold=True, space_after=0)

prs.save(OUT)
print("salvo:", OUT, "| slides:", len(prs.slides._sldIdLst))
