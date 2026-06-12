# -*- coding: utf-8 -*-
"""Gera a apresentação (PPTX) reequilibrada para o Plano B — Equipe 4."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ---------- paleta sóbria ----------
INK    = RGBColor(0x2B, 0x2B, 0x2B)
MUTE   = RGBColor(0x6B, 0x72, 0x76)
ACCENT = RGBColor(0x1F, 0x4E, 0x5F)   # azul-petróleo
ACCENT2= RGBColor(0x3C, 0x7A, 0x8C)
HAIR   = RGBColor(0xD9, 0xDE, 0xE0)
PANEL  = RGBColor(0xF4, 0xF6, 0xF7)
PAPER  = RGBColor(0xFF, 0xFF, 0xFF)
LTBLUE = RGBColor(0xCD, 0xDE, 0xE3)
S_GREY = RGBColor(0x9A, 0xA0, 0xA4)
S_AMBER= RGBColor(0xB0, 0x7A, 0x2A)
S_GREEN= RGBColor(0x2E, 0x6B, 0x4F)

SERIF = "Cambria"
SANS  = "Calibri"

BASE = "/home/user/reds-dm1-otimizador-hly/plano_b_pso_diabetes"
CONV_PNG = BASE + "/resultados_convergencia_pso.png"
IMP_PNG  = BASE + "/resultados_importancia.png"
OUT = BASE + "/Apresentacao_Equipe4_PlanoB.pptx"

ARROW = "→"  # →

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
ML = Inches(0.92)
CW = SW - ML - Inches(0.92)

def _set_font(run, size=18, color=INK, bold=False, italic=False, font=SANS):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = font
    f.color.rgb = color

def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf

def para(tf, text, size=18, color=INK, bold=False, italic=False, font=SANS,
         align=PP_ALIGN.LEFT, space_after=6, space_before=0, level=0, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.level = level
    p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    r = p.add_run(); r.text = text
    _set_font(r, size, color, bold, italic, font)
    return p, r

def rect(slide, l, t, w, h, fill=PANEL, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp

def hline(slide, l, t, w, color=ACCENT, weight=1.5):
    ln = slide.shapes.add_connector(2, l, t, l + w, t)
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln

def chrome(slide, kicker, title, idx, total=13):
    _, ktf = box(slide, ML, Inches(0.45), CW, Inches(0.3))
    para(ktf, kicker.upper(), size=11, color=ACCENT2, bold=True, font=SANS, space_after=0, first=True)
    _, ttf = box(slide, ML, Inches(0.74), CW, Inches(0.7))
    para(ttf, title, size=27, color=INK, bold=True, font=SERIF, space_after=0, first=True)
    hline(slide, ML, Inches(1.52), CW, color=ACCENT, weight=1.5)
    _, ftf = box(slide, ML, SH - Inches(0.5), Inches(8), Inches(0.3))
    para(ftf, "Projeto REDS-DM1  ·  Equipe 4  ·  Computação Natural",
         size=9, color=MUTE, font=SANS, space_after=0, first=True)
    _, ntf = box(slide, SW - Inches(1.8), SH - Inches(0.5), Inches(0.88), Inches(0.3))
    para(ntf, f"{idx:02d} / {total:02d}", size=9, color=MUTE, font=SANS,
         align=PP_ALIGN.RIGHT, space_after=0, first=True)

def newslide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=PAPER)
    return s

def bullets(tf, items, size=16, gap=8, color=INK, lead=ACCENT):
    for i, it in enumerate(items):
        txt, bold = (it if isinstance(it, tuple) else (it, False))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        r0 = p.add_run(); r0.text = "—  "; _set_font(r0, size, lead, True, False, SANS)
        r = p.add_run(); r.text = txt; _set_font(r, size, color, bold, False, SANS)

# ===================== CAPA =====================
s = newslide()
rect(s, 0, 0, Inches(0.22), SH, fill=ACCENT)
_, ktf = box(s, ML, Inches(1.55), CW, Inches(0.4))
para(ktf, "COMPUTAÇÃO NATURAL  ·  EQUIPE 4", size=13, color=ACCENT2, bold=True, font=SANS, first=True)
_, ttf = box(s, ML, Inches(2.05), CW, Inches(2.0))
para(ttf, "Projeto REDS-DM1", size=40, color=INK, bold=True, font=SERIF, first=True, space_after=4)
para(ttf, "Otimizador Bioinspirado para Gestão de Saúde em Diabetes",
     size=24, color=ACCENT, font=SERIF, space_after=0)
hline(s, ML, Inches(3.95), Inches(4.2), color=ACCENT, weight=2)
_, stf = box(s, ML, Inches(4.2), CW, Inches(1.2))
para(stf, "Entrega: PSO validado sobre coorte de diabéticos  (Plano B)",
     size=18, color=INK, bold=True, font=SANS, first=True, space_after=4)
para(stf, "Inteligência de Enxames  ·  PSO  ·  FSS  ·  proxy de HLY",
     size=14, color=MUTE, font=SANS, space_after=0)
_, dtf = box(s, ML, SH - Inches(0.7), CW, Inches(0.4))
para(dtf, "REDS-PE  —  Rede de Atenção à Saúde  ·  Pernambuco", size=11, color=MUTE, font=SANS, first=True)

# ===================== 1 AGENDA =====================
s = newslide()
chrome(s, "Roteiro", "Agenda", 1)
top = Inches(1.95)
items = [
    ("01", "De onde viemos", "Trajetória em 3 fases e redução de escopo"),
    ("02", "O Plano B e seu objetivo", "Objetivo adaptado e a variável Wellness"),
    ("03", "Como funciona", "Passo a passo e por que PSO (enxame)"),
    ("04", "Validação e ganho", "O ótimo conhecido e o perfil do indivíduo ideal"),
    ("05", "Fechamento", "BFSS de relance, limitações e próximos passos"),
]
rowh = Inches(0.92); gap = Inches(0.05)
for i,(n,t,d) in enumerate(items):
    y = top + i*(rowh+gap); hi = i in (1,2,3)
    rect(s, ML, y, CW, rowh, fill=(PANEL if hi else PAPER), line=HAIR, line_w=0.75)
    if hi: rect(s, ML, y, Inches(0.08), rowh, fill=ACCENT)
    _, ntf = box(s, ML+Inches(0.3), y+Inches(0.18), Inches(0.9), Inches(0.6))
    para(ntf, n, size=24, color=ACCENT2, bold=True, font=SERIF, first=True)
    _, ttf = box(s, ML+Inches(1.35), y+Inches(0.13), CW-Inches(1.6), Inches(0.7))
    para(ttf, t, size=17, color=INK, bold=True, font=SANS, first=True, space_after=1)
    para(ttf, d, size=12.5, color=MUTE, font=SANS, space_after=0)

# ===================== 2 SÍNTESE =====================
s = newslide()
chrome(s, "Síntese executiva", "O que é, qual o problema, onde chegamos", 2)
colw = (CW - Inches(0.6)) / 3
cols = [
    ("Contexto", [
        "Disciplina de Computação Natural.",
        "Sistema REDS-PE (Rede de Atenção à Saúde).",
        "Alvo: pacientes com diabetes em Pernambuco.",
        "Financiamento: SUS / políticas públicas.",
    ]),
    ("Proposta", [
        "Usar Inteligência de Enxames (PSO).",
        "Achar o perfil do 'indivíduo ideal' entre diabéticos.",
        "Horizonte: otimizar a alocação de recursos maximizando HLY.",
    ]),
    ("Resultado", [
        "PSO implementado do zero (NumPy).",
        "Validado contra a solução analítica: cosseno = 1,000.",
        "Perfil clínico interpretável extraído dos pesos.",
    ]),
]
ctop = Inches(1.95)
for i,(h,its) in enumerate(cols):
    x = ML + i*(colw+Inches(0.3))
    rect(s, x, ctop, colw, Inches(4.4), fill=PANEL, line=HAIR, line_w=0.75)
    barcol = ACCENT if i<2 else S_GREEN
    rect(s, x, ctop, colw, Inches(0.12), fill=barcol)
    _, tf = box(s, x+Inches(0.28), ctop+Inches(0.4), colw-Inches(0.56), Inches(3.8))
    para(tf, h, size=18, color=barcol, bold=True, font=SERIF, first=True, space_after=10)
    bullets(tf, its, size=14, gap=10, lead=barcol)

# ===================== 3 TRAJETÓRIA =====================
s = newslide()
chrome(s, "De onde viemos", "Trajetória em três fases", 3)
_, sub = box(s, ML, Inches(1.65), CW, Inches(0.4))
para(sub, "A redução de escopo foi decisão técnica, não recuo.", size=14, color=MUTE, italic=True, font=SANS, first=True)
phases = [
    ("FASE 1", "130-US + CTGAN", "Abandonada", S_GREY,
     ["População errada (97% DM2)", "Sem eixo temporal", "Sem gabarito de validação"]),
    ("FASE 2", "Marcadores plantados", "Pausada", S_AMBER,
     ["Base sintética com gabarito", "Prova algorítmica do otimizador", "BFSS validado aqui"]),
    ("FASE 3", "Plano B — PSO / coorte bimodal", "Entregue", S_GREEN,
     ["Foco da apresentação de hoje", "PSO validado (cosseno 1,000)", "Perfil ideal interpretável"]),
]
cw3 = (CW - Inches(1.0)) / 3
ytop = Inches(2.35); yline = ytop
hline(s, ML+Inches(0.3), yline, CW-Inches(0.6), color=HAIR, weight=2)
for i,(ph,name,st,col,pts) in enumerate(phases):
    x = ML + i*(cw3+Inches(0.5))
    rect(s, x+cw3/2-Inches(0.12), yline-Inches(0.12), Inches(0.24), Inches(0.24), fill=col, shape=MSO_SHAPE.OVAL)
    card_t = ytop + Inches(0.45)
    rect(s, x, card_t, cw3, Inches(3.4), fill=PAPER, line=HAIR, line_w=0.75)
    rect(s, x, card_t, cw3, Inches(0.1), fill=col)
    _, tf = box(s, x+Inches(0.25), card_t+Inches(0.3), cw3-Inches(0.5), Inches(3.0))
    para(tf, ph, size=11, color=MUTE, bold=True, font=SANS, first=True, space_after=2)
    para(tf, name, size=16.5, color=INK, bold=True, font=SERIF, space_after=6)
    para(tf, st.upper(), size=11, color=col, bold=True, font=SANS, space_after=10)
    bullets(tf, pts, size=12.5, gap=7, lead=col)

# ===================== 4 REDUÇÃO DE ESCOPO =====================
s = newslide()
chrome(s, "Transparência", "Redução de escopo — status de cada componente", 4)
rows = [
    ("130-US + CTGAN", "Abandonado", S_GREY),
    ("PSO sobre coorte bimodal (Plano B)", "Entregue", S_GREEN),
    ("Meta-FSS (hiperparâmetros do PSO)", "Entregue", S_GREEN),
    ("Visualização HTML do PSO (3 abas)", "Entregue", S_GREEN),
    ("Banco sintético de marcadores", "Concluído", S_GREEN),
    ("BFSS (seleção de variáveis)", "Validado", S_GREEN),
    ("MOPSO + NSGA-II (frente de Pareto)", "Planejado", S_AMBER),
]
ytop = Inches(1.85); rh = Inches(0.62); gap = Inches(0.07)
for i,(comp,st,col) in enumerate(rows):
    y = ytop + i*(rh+gap); bold = "Plano B" in comp
    rect(s, ML, y, CW, rh, fill=(PANEL if bold else PAPER), line=HAIR, line_w=0.75)
    rect(s, ML, y, Inches(0.1), rh, fill=col)
    _, tf = box(s, ML+Inches(0.35), y+Inches(0.14), CW-Inches(3.0), Inches(0.4))
    para(tf, comp, size=15.5, color=INK, bold=bold, font=SANS, first=True)
    bw = Inches(1.9); bx = ML + CW - bw - Inches(0.2)
    rect(s, bx, y+Inches(0.13), bw, Inches(0.36), fill=None, line=col, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _, btf = box(s, bx, y+Inches(0.16), bw, Inches(0.34))
    para(btf, st.upper(), size=11, color=col, bold=True, font=SANS, align=PP_ALIGN.CENTER, first=True)

# ===================== 5 OBJETIVO =====================
s = newslide()
chrome(s, "O Plano B", f"Objetivo: original  {ARROW}  adaptado", 5)
colw = (CW - Inches(1.1)) / 2; ytop = Inches(1.95); ch = Inches(3.7)
rect(s, ML, ytop, colw, ch, fill=PAPER, line=HAIR, line_w=0.75)
rect(s, ML, ytop, colw, Inches(0.12), fill=S_GREY)
_, tf = box(s, ML+Inches(0.3), ytop+Inches(0.35), colw-Inches(0.6), ch-Inches(0.7))
para(tf, "ORIGINAL — PLANO A", size=12, color=MUTE, bold=True, font=SANS, first=True, space_after=8)
para(tf, "Otimizador multiobjetivo completo", size=17, color=INK, bold=True, font=SERIF, space_after=10)
bullets(tf, [
    "Três objetivos concorrentes (frente de Pareto):",
    "Maximizar HLY (Anos de Vida Saudáveis).",
    "Minimizar o custo da política pública.",
    "Maximizar a equidade regional entre municípios.",
    "Baselines: MOPSO + NSGA-II.",
], size=13.5, gap=8, lead=S_GREY)
ax = ML + colw + Inches(0.18)
_, atf = box(s, ax, ytop+ch/2-Inches(0.45), Inches(0.7), Inches(0.8))
para(atf, ARROW, size=40, color=ACCENT, bold=True, font=SANS, align=PP_ALIGN.CENTER, first=True)
x2 = ML + colw + Inches(1.0)
rect(s, x2, ytop, colw, ch, fill=PANEL, line=ACCENT, line_w=1.25)
rect(s, x2, ytop, colw, Inches(0.12), fill=ACCENT)
_, tf2 = box(s, x2+Inches(0.3), ytop+Inches(0.35), colw-Inches(0.6), ch-Inches(0.7))
para(tf2, "ADAPTADO — PLANO B  (entregável real)", size=12, color=ACCENT, bold=True, font=SANS, first=True, space_after=8)
para(tf2, "PSO acha o perfil do indivíduo ideal", size=17, color=INK, bold=True, font=SERIF, space_after=10)
bullets(tf2, [
    "PSO otimiza os pesos de uma regressão logística.",
    "Variável-alvo: Wellness (proxy de HLY).",
    "Validar contra a solução analítica do scikit-learn.",
    "Extrair o perfil clínico interpretável dos pesos.",
], size=13.5, gap=8, lead=ACCENT)
_, ntf = box(s, ML, ytop+ch+Inches(0.18), CW, Inches(0.5))
para(ntf, "Honestidade metodológica: a base é testbed; o entregável é o otimizador. Coorte única (China, 2012), que não distingue T1D/T2D.",
     size=12, color=MUTE, italic=True, font=SANS, first=True)

# ===================== 6 WELLNESS =====================
s = newslide()
chrome(s, "O Plano B", "A variável dependente: Wellness", 6)
lw = CW*0.56
_, tf = box(s, ML, Inches(1.95), lw, Inches(4.2))
para(tf, "Por que não um rótulo clínico direto?", size=17, color=INK, bold=True, font=SERIF, first=True, space_after=10)
bullets(tf, [
    "A base não possui rótulo 'saudável vs. doente'.",
    "É um snapshot transversal, sem trajetória longitudinal.",
    "Wellness é construído a partir de variáveis objetivas da coorte.",
], size=14.5, gap=10, lead=ACCENT)
para(tf, "Conexão com HLY", size=15, color=ACCENT, bold=True, font=SANS, space_after=6, space_before=8)
para(tf, "Ausência de complicações + controle metabólico = pilares que sustentam os Anos de Vida Saudáveis.",
     size=14, color=INK, font=SANS, space_after=0)
fy = Inches(4.05)
rect(s, ML, fy, lw, Inches(0.95), fill=PANEL, line=HAIR, line_w=0.75)
_, ftf = box(s, ML+Inches(0.3), fy+Inches(0.12), lw-Inches(0.6), Inches(0.75))
para(ftf, "Wellness  =  terço superior de", size=13, color=MUTE, font=SANS, first=True, space_after=2)
para(ftf, "z(idade) − z(HbA1c) − z(n.º de comorbidades)", size=18, color=ACCENT, bold=True, font=SERIF, space_after=0)
rx = ML + lw + Inches(0.5); rw = CW - lw - Inches(0.5); gy = Inches(1.95)
rect(s, rx, gy, rw, Inches(3.0), fill=PAPER, line=HAIR, line_w=0.75)
rect(s, rx, gy, rw, Inches(0.55), fill=ACCENT)
_, htf = box(s, rx, gy+Inches(0.12), rw, Inches(0.4))
para(htf, "Grupos  (filtro DM == 1)", size=14, color=PAPER, bold=True, font=SANS, align=PP_ALIGN.CENTER, first=True)
grp = [("Indivíduos ideais", "119", "terço superior"), ("Demais diabéticos", "238", ""), ("Total", "357", "")]
yy = gy + Inches(0.72)
for i,(g,n,note) in enumerate(grp):
    if i==2:
        hline(s, rx+Inches(0.2), yy, rw-Inches(0.4), color=HAIR, weight=1); yy += Inches(0.08)
    _, rtf = box(s, rx+Inches(0.3), yy, rw-Inches(1.4), Inches(0.5))
    para(rtf, g, size=15, color=INK, bold=(i==2 or i==0), font=SANS, first=True, space_after=0)
    if note: para(rtf, note, size=10.5, color=MUTE, italic=True, font=SANS, space_after=0)
    _, vtf = box(s, rx+rw-Inches(1.2), yy-Inches(0.02), Inches(1.0), Inches(0.5))
    para(vtf, n, size=20, color=ACCENT, bold=True, font=SERIF, align=PP_ALIGN.RIGHT, first=True)
    yy += Inches(0.72)

# ===================== 7 PASSO A PASSO =====================
s = newslide()
chrome(s, "Como funciona", "Passo a passo da solução (Plano B)", 7)
_, sub = box(s, ML, Inches(1.62), CW, Inches(0.35))
para(sub, "Da base bruta ao perfil do indivíduo ideal — 9 etapas reprodutíveis.", size=13.5, color=MUTE, italic=True, font=SANS, first=True)
groups = [
    ("PREPARAÇÃO", ACCENT2, [
        ("1", "Carregar & filtrar", "CSV 5.922×190 → DM==1 → 357"),
        ("2", "Construir Wellness", "rótulo 0/1"),
        ("3", "Selecionar 26 features", "imputar mediana; z-score"),
    ]),
    ("OTIMIZAÇÃO", ACCENT, [
        ("4", "Função de fitness", "−(log-loss + L2)"),
        ("5", "Executar PSO", "27D; w=0,7; c1=c2=1,5"),
        ("6", "Validar", "cosseno PSO × sklearn = 1,000"),
    ]),
    ("INTERPRETAÇÃO", S_GREEN, [
        ("7", "Meta-FSS", "hiperparâm.; AUC cross-val"),
        ("8", "Perfil do ideal", "rank por |peso|; sinal = direção"),
        ("9", "Visualizar", "HTML, 3 abas, enxame animado"),
    ]),
]
gw = (CW - Inches(0.8))/3; gy = Inches(2.15)
for i,(gh,col,steps) in enumerate(groups):
    x = ML + i*(gw+Inches(0.4))
    _, ghtf = box(s, x, gy, gw, Inches(0.35))
    para(ghtf, gh, size=12, color=col, bold=True, font=SANS, align=PP_ALIGN.CENTER, first=True)
    sy = gy + Inches(0.5)
    for (n,t,d) in steps:
        rect(s, x, sy, gw, Inches(1.15), fill=PANEL, line=HAIR, line_w=0.75)
        rect(s, x, sy, Inches(0.12), Inches(1.15), fill=col)
        rect(s, x+Inches(0.28), sy+Inches(0.32), Inches(0.5), Inches(0.5), fill=col, shape=MSO_SHAPE.OVAL)
        _, ntf = box(s, x+Inches(0.28), sy+Inches(0.42), Inches(0.5), Inches(0.4))
        para(ntf, n, size=16, color=PAPER, bold=True, font=SANS, align=PP_ALIGN.CENTER, first=True)
        _, stf = box(s, x+Inches(0.95), sy+Inches(0.2), gw-Inches(1.1), Inches(0.85))
        para(stf, t, size=14.5, color=INK, bold=True, font=SANS, first=True, space_after=2)
        para(stf, d, size=11.5, color=MUTE, font=SANS, space_after=0)
        sy += Inches(1.25)

# ===================== 8 POR QUE PSO =====================
s = newslide()
chrome(s, "Como funciona", "Por que PSO (enxame) e não algoritmo evolucionário", 8)
quad = [
    ("Terreno do problema", ACCENT,
     "Os pesos vivem em espaço contínuo de alta dimensão (27D) e a perda é convexa. "
     "O PSO desliza no contínuo com momento e converge rápido para o ótimo. "
     "GAs exigem crossover/mutação, menos intuitivos para pesos reais."),
    ("Alinhamento com o enunciado", ACCENT2,
     "A disciplina avalia Inteligência de Enxames. PSO é o algoritmo canônico, "
     "e a família FSS é o fio condutor do projeto."),
    ("Interpretabilidade", ACCENT2,
     "O PSO entrega direto um vetor de pesos: magnitude = importância, "
     "sinal = direção clínica. Um genoma binário de GA obscureceria isso."),
    ("Validação elegante", S_GREEN,
     "Em problema convexo existe ótimo conhecido — dá para provar que o "
     "otimizador chegou nele (cosseno = 1,000). Demonstração limpa e didática."),
]
qw = (CW - Inches(0.5))/2; qh = Inches(2.15); qy0 = Inches(1.9)
for i,(h,col,body) in enumerate(quad):
    r = i//2; c = i%2
    x = ML + c*(qw+Inches(0.5)); y = qy0 + r*(qh+Inches(0.3))
    rect(s, x, y, qw, qh, fill=PANEL, line=HAIR, line_w=0.75)
    rect(s, x, y, Inches(0.12), qh, fill=col)
    _, tf = box(s, x+Inches(0.35), y+Inches(0.25), qw-Inches(0.65), qh-Inches(0.5))
    para(tf, h, size=16.5, color=col, bold=True, font=SERIF, first=True, space_after=7)
    para(tf, body, size=13.5, color=INK, font=SANS, space_after=0)

# ===================== 9 VALIDAÇÃO =====================
s = newslide()
chrome(s, "Validação e ganho", "O otimizador acha o ótimo conhecido", 9)
lw = CW*0.42
metrics = [("AUC treino", "0,820", MUTE), ("AUC teste", "0,664", INK), ("AUC sklearn (baseline)", "0,664", INK)]
my = Inches(2.0)
for (lbl,val,col) in metrics:
    rect(s, ML, my, lw, Inches(0.7), fill=PAPER, line=HAIR, line_w=0.75)
    _, ltf = box(s, ML+Inches(0.25), my+Inches(0.18), lw-Inches(1.5), Inches(0.4))
    para(ltf, lbl, size=14, color=col, font=SANS, first=True)
    _, vtf = box(s, ML+lw-Inches(1.4), my+Inches(0.09), Inches(1.2), Inches(0.55))
    para(vtf, val, size=22, color=col, bold=True, font=SERIF, align=PP_ALIGN.RIGHT, first=True)
    my += Inches(0.82)
cy = my + Inches(0.1)
rect(s, ML, cy, lw, Inches(1.35), fill=ACCENT)
_, ctf = box(s, ML+Inches(0.25), cy+Inches(0.2), lw-Inches(0.5), Inches(1.0))
para(ctf, "Cosseno(pesos PSO, sklearn)", size=13, color=LTBLUE, bold=True, font=SANS, first=True, space_after=2)
para(ctf, "1,000", size=40, color=PAPER, bold=True, font=SERIF, space_after=0)
rx = ML + lw + Inches(0.5); rw = CW - lw - Inches(0.5)
rect(s, rx, Inches(1.95), rw, Inches(3.6), fill=PAPER, line=HAIR, line_w=0.75)
s.shapes.add_picture(CONV_PNG, rx+Inches(0.15), Inches(2.12), width=rw-Inches(0.3))
_, ntf = box(s, ML, Inches(5.95), CW, Inches(0.8))
para(ntf, "O AUC mede a dificuldade do problema; o cosseno prova o otimizador. Para uma demonstração de "
          "validação, o PSO num problema convexo é limpo e didático: encontra EXATAMENTE a mesma solução que o "
          "otimizador analítico.", size=12.5, color=MUTE, italic=True, font=SANS, first=True)

# ===================== 10 GANHO / PERFIL =====================
s = newslide()
chrome(s, "Validação e ganho", "O ganho: o perfil do indivíduo ideal", 10)
rw = CW*0.5; rx = ML + CW - rw
rect(s, rx, Inches(1.9), rw, Inches(4.6), fill=PAPER, line=HAIR, line_w=0.75)
s.shapes.add_picture(IMP_PNG, rx+Inches(0.2), Inches(2.05), height=Inches(4.3))
lw = CW - rw - Inches(0.5)
_, tf = box(s, ML, Inches(1.95), lw, Inches(2.2))
para(tf, "Top características (peso do PSO)", size=15, color=ACCENT, bold=True, font=SANS, first=True, space_after=8)
bullets(tf, [
    "CP2h +0,338  ·  BMI −0,304  ·  BUN +0,240",
    "waist1 −0,232  ·  ISIGutt −0,212  ·  WHR −0,193",
    "SCRE +0,162  ·  ALT −0,134  ·  FatherDM −0,119",
], size=13.5, gap=7, lead=ACCENT)
py = Inches(3.55)
rect(s, ML, py, lw, Inches(1.55), fill=PANEL, line=ACCENT, line_w=1.0)
_, ptf = box(s, ML+Inches(0.25), py+Inches(0.18), lw-Inches(0.5), Inches(1.2))
para(ptf, "Perfil-síntese do diabético ideal", size=14, color=ACCENT, bold=True, font=SANS, first=True, space_after=5)
para(ptf, "Mais magro, com menor resistência à insulina, função hepática preservada (ALT) e sem histórico familiar de DM.",
     size=13.5, color=INK, font=SANS, space_after=0)
my = Inches(5.3)
_, mtf = box(s, ML, my, lw, Inches(1.2))
para(mtf, "Médias (ideais vs demais):", size=13, color=MUTE, bold=True, font=SANS, first=True, space_after=4)
para(mtf, "BMI 21,4 vs 24,8   ·   ISIGutt 57,3 vs 104,8   ·   ALT 33,4 vs 45,7", size=13, color=INK, font=SANS, space_after=0)

# ===================== 11 BFSS DE RELANCE =====================
s = newslide()
chrome(s, "Fechamento", "BFSS, de relance — o plano principal segue vivo", 11)
_, tf = box(s, ML, Inches(2.0), CW*0.58, Inches(4.0))
para(tf, "Validação por marcadores plantados", size=17, color=INK, bold=True, font=SERIF, first=True, space_after=10)
bullets(tf, [
    "Base sintética com gabarito (verdade-base plantada).",
    "Prova ALGORÍTMICA de que o otimizador recupera o sinal — não é alegação clínica.",
    "Mantém a família de enxame (FSS) como fio condutor do projeto.",
], size=14.5, gap=10, lead=ACCENT)
para(tf, f"Próximo: MOPSO discreto + NSGA-II para a frente de Pareto.",
     size=14, color=MUTE, italic=True, font=SANS, space_after=0, space_before=6)
rx = ML + CW*0.62; rw = CW*0.38
rect(s, rx, Inches(2.0), rw, Inches(1.9), fill=PANEL, line=HAIR, line_w=0.75)
rect(s, rx, Inches(2.0), rw, Inches(0.1), fill=S_GREEN)
_, stf = box(s, rx+Inches(0.3), Inches(2.25), rw-Inches(0.6), Inches(1.5))
para(stf, "STATUS", size=11, color=MUTE, bold=True, font=SANS, first=True, space_after=4)
para(stf, "BFSS validado", size=18, color=S_GREEN, bold=True, font=SERIF, space_after=6)
para(stf, "Recuperou 8 de 9 marcadores plantados; o subconjunto selecionado eleva o R².",
     size=13, color=INK, font=SANS, space_after=0)

# ===================== 12 LIMITAÇÕES =====================
s = newslide()
chrome(s, "Fechamento", "Limitações declaradas", 12)
lims = [
    "A base não distingue T1D/T2D → 'diabetes tipo não especificado'; é testbed, não evidência clínica.",
    "Coorte única e transversal (China, 2012) → fala de associação, não de causalidade.",
    "O rótulo 'ideal' é uma escolha de projeto defensável, não a única.",
    "As 26 preditoras foram curadas à mão → podem ter omitido pistas relevantes.",
    "Imputação pela mediana 'achata' a variação; o corte do terço (33%) é limiar arbitrário.",
    "L2 (λ=0,05) calibrado empiricamente, não por cross-validation formal.",
]
colw = (CW - Inches(0.5))/2; ytop = Inches(2.0)
for i,t in enumerate(lims):
    r = i//2; c = i%2
    x = ML + c*(colw+Inches(0.5)); y = ytop + r*Inches(1.45)
    rect(s, x, y, colw, Inches(1.25), fill=PANEL, line=HAIR, line_w=0.75)
    rect(s, x, y, Inches(0.1), Inches(1.25), fill=ACCENT2)
    _, tf = box(s, x+Inches(0.3), y+Inches(0.22), colw-Inches(0.55), Inches(0.95))
    para(tf, t, size=13.5, color=INK, font=SANS, first=True)

# ===================== 13 CONCLUSÃO =====================
s = newslide()
chrome(s, "Fechamento", "Próximos passos e conclusão", 13)
stages = ["PSO (Plano B)", "BFSS", "Simulador HLY", "MOPSO + NSGA-II"]
pw = (CW - Inches(1.2))/4; py = Inches(2.1)
for i,st in enumerate(stages):
    x = ML + i*(pw+Inches(0.4))
    col = S_GREEN if i<2 else (S_AMBER if i==3 else ACCENT2)
    rect(s, x, py, pw, Inches(0.85), fill=PANEL, line=col, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _, tf = box(s, x, py+Inches(0.26), pw, Inches(0.5))
    para(tf, st, size=14, color=INK, bold=True, font=SANS, align=PP_ALIGN.CENTER, first=True)
    if i < 3:
        _, atf = box(s, x+pw+Inches(0.02), py+Inches(0.2), Inches(0.36), Inches(0.5))
        para(atf, ARROW, size=20, color=MUTE, bold=True, font=SANS, align=PP_ALIGN.CENTER, first=True)
_, tf = box(s, ML, Inches(3.4), CW, Inches(1.2))
para(tf, "Próximo passo — frente de Pareto", size=16, color=ACCENT, bold=True, font=SERIF, first=True, space_after=6)
para(tf, "MOPSO discreto + NSGA-II sobre ~540 políticas, com gabarito computável por enumeração "
         "(HLY × Custo × Equidade).", size=14.5, color=INK, font=SANS, space_after=0)
fy = Inches(4.7)
rect(s, ML, fy, CW, Inches(1.4), fill=ACCENT)
_, ftf = box(s, ML+Inches(0.4), fy+Inches(0.28), CW-Inches(0.8), Inches(1.0))
para(ftf, "Conclusão", size=13, color=LTBLUE, bold=True, font=SANS, first=True, space_after=4)
para(ftf, "PSO validado (cosseno 1,000) + perfil clínico interpretável = o otimizador funciona. "
          "O Plano A continua como horizonte multiobjetivo.", size=17, color=PAPER, bold=True, font=SERIF, space_after=0)

prs.save(OUT)
print("salvo:", OUT, "| slides:", len(prs.slides._sldIdLst))
